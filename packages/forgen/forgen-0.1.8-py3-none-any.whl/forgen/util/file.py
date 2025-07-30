import base64
import hashlib
import io
import json
import os
import secrets
import string
from io import BytesIO

import chardet
import docx2txt
import pandas as pd
import pytesseract
import requests
from PIL import Image
from bs4 import BeautifulSoup
from doc2docx import convert
from docx import Document
from pdf2image import convert_from_bytes
from pptx import Presentation
from pypdf import PdfReader
from pytesseract import image_to_pdf_or_hocr
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from striprtf.striprtf import rtf_to_text


def perform_general_ocr(file_stream):
    file_stream.seek(0)
    images = convert_from_bytes(file_stream.read())

    ocr_results = []

    for page_number, image in enumerate(images, start=1):
        hocr = image_to_pdf_or_hocr(image, extension='hocr', lang='eng')
        soup = BeautifulSoup(hocr, 'html.parser')

        for line in soup.find_all('span', class_='ocr_line'):
            line_text = line.get_text(strip=False).replace('\n', ' ').strip()
            bbox = line['title'].split(';')[0].split(' ')[1:]
            if len(bbox) == 4:
                x0, y0, x1, y1 = bbox
                ocr_results.append({
                    'page': page_number,
                    'text': line_text,
                    'bbox': {'x0': x0, 'y0': y0, 'x1': x1, 'y1': y1}
                })

    return ocr_results

def random_string(length=10):
    """Generate a random string of letters and digits."""
    characters = string.ascii_letters + string.digits
    return ''.join(secrets.choice(characters) for i in range(length))

def download_pdf_from_url(url, output_path):
    response = requests.get(url)
    if response.status_code == 200:
        with open(output_path, 'wb') as file:
            file.write(response.content)
        print(f"PDF successfully downloaded to {output_path}")
    else:
        print(f"Failed to download PDF. Status code: {response.status_code} : {url}")


def extract_text_from_rtf(file_path):
    with open(file_path, 'r') as file:
        rtf_text = file.read()
    the_text = rtf_to_text(rtf_text)
    print(f"failed to extract text from docx: {str(the_text)}")
    return the_text


def extract_text_from_doc(file_obj, file_path):
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_name = hash_file_or_string(file_obj) + '-tempfile.docx'
            docx_path = os.path.join(temp_dir, temp_name)
            with open(docx_path, 'wb') as f:
                f.write(file_obj.read())
            docx_path = os.path.join(temp_dir, temp_name)
            convert(file_path, docx_path)
            extracted_text = docx2txt.process(docx_path)
    except Exception as e:
        print(f"failed to extract text from docx: {str(e)}")
        try:
            return extract_text_from_rtf(file_path)
        except Exception as e:
            print(f"failed to extract text from docx: {str(e)}")
            return f"failed to extract text: {str(e)}"
    return extracted_text


def convert_and_extract_text_from_doc(file_obj):
    """
    Converts a .doc file, given as a file object, to .docx format and extracts its text.
    :param file_obj: file object, file object of the .doc file to be converted.
    :return: str, extracted text from the converted .docx file.
    """
    # Create a temporary directory to save the doc file and the converted docx file
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save the original .doc file from the file object to the temporary directory
            temp_name = hash_file_or_string(file_obj) + '-tempfile.docx'
            doc_path = os.path.join(temp_dir, temp_name)
            with open(doc_path, 'wb') as f:
                f.write(file_obj.read())
            docx_path = os.path.join(temp_dir, temp_name)
            convert(doc_path, docx_path)
            extracted_text = docx2txt.process(docx_path)
    except:
        return "failed to extract text from docx"
    return extracted_text


def extract_text_from_docx(file):
    text = docx2txt.process(file)
    return text


def extract_text_from_pptx(file):
    presentation = Presentation(file)
    extracted_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        extracted_text.append(run.text)

    # Join the extracted text into a single string
    extracted_text = ' '.join(extracted_text)
    return extracted_text


def load_pdf_to_file_object(file_path):
    with open(file_path, 'rb') as f:
        pdf_data = f.read()
    return io.BytesIO(pdf_data)


def load_docx_to_file_object(file_path):
    doc = Document(file_path)
    file_object = io.BytesIO()
    doc.save(file_object)
    file_object.seek(0)
    return file_object


def fetch_image_and_extract_text(image_url):
    response = requests.get(image_url)
    image = Image.open(BytesIO(response.content))
    extracted_text = pytesseract.image_to_string(image)
    return extracted_text


def fetch_image_and_convert_to_pdf(image_url):
    try:
        response = requests.get(image_url)
        response.raise_for_status()
        image = Image.open(BytesIO(response.content))
        temp_image_path = f"temp_image-{random_string(10)}.png"
        image.save(temp_image_path)
        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=letter)
        c.drawImage(temp_image_path, 0, 0, width=letter[0], height=letter[1])
        c.save()
        pdf_buffer.seek(0)
        os.remove(temp_image_path)
    except (requests.RequestException, IOError, Exception) as e:
        print(f"Error during image download or conversion to PDF: {e}")
        return None
    return pdf_buffer


def fetch_image_and_extract_text_using_ocr(image_url):
    pdf_file = fetch_image_and_convert_to_pdf(image_url)
    if pdf_file is None:
        return None
    ocr_result = extract_text_from_pdf(pdf_file)
    return ocr_result


def get_text_from_file(file, filename):
    print(f"get_text_from_file filename: {filename}")
    if file:
        file_ext = filename.split(".")[-1].lower()
        print(f"get_text_from_file file_ext: {file_ext}")
        if file_ext in ["txt", "xml"]:
            encoding = determine_encoding(file)
            print(f"get_text_from_file encoding: {encoding}")
            if encoding is None:
                encoding = "utf-8"
            text = file.read().decode(encoding, errors="replace")
            print(f"get_text_from_file txt: {text}")
        elif file_ext in ["docx"]:
            text = extract_text_from_docx(file)
        elif file_ext in ["doc"]:
            text = convert_and_extract_text_from_doc(file)
        elif file_ext in ["pptx", "ppt"]:
            text = extract_text_from_pptx(file)
        elif file_ext in ["pdf"]:
            output_obj = extract_text_from_pdf(file)
            try:
                output_obj = json.loads(output_obj)
                print(f"get_text_from_file json loads success: {str(output_obj)[:100]}")
                sorted_data = sorted(output_obj, key=lambda x: (x['page'], x['column'], x['line_number']))
                document_text = "\n".join([item['text'] for item in sorted_data if item['text']])
                return document_text, output_obj
            except Exception as e:
                print(f"get_text_from_file extract text from pdf error when json loading")
                pass
            if isinstance(output_obj, list):
                return str(output_obj), output_obj
            text = str(output_obj)
        elif file_ext in ["csv"]:
            df = pd.read_csv(file)
            text = df.to_string()
        elif file_ext in ["xls", "xlsx"]:
            df = pd.read_excel(file)
            text = df.to_string()
        else:
            raise ValueError("Unsupported file type")
        return text, None
    else:
        raise ValueError("No file provided")


def determine_encoding(file_object):
    raw_data = file_object.read()
    result = chardet.detect(raw_data)
    encoding = result["encoding"]
    file_object.seek(0)
    return encoding


def pdf_has_text(file_stream):
    try:
        reader = PdfReader(file_stream)
        for page in reader.pages:
            if page.extract_text():
                return True
        return False
    except Exception as e:
        return False


# from pdf2image import convert_from_path
import tempfile

# def ocr_pdf(file):
#     # Use a with statement to ensure temp files are cleaned up after use
#     with tempfile.NamedTemporaryFile(suffix=".pdf") as temp_pdf:
#         # Write the uploaded file's content to the temp file
#         temp_pdf.write(file.read())
#         temp_pdf.seek(0)
#
#         # Convert the PDF to a list of image objects
#         # images = convert_from_path(temp_pdf.name, poppler_path='/usr/bin/poppler')
#         images = convert_from_path(temp_pdf.name)
#
#         # Use Tesseract to do OCR on the images
#         text = ""
#         for image in images:
#             text += pytesseract.image_to_string(image)
#
#         return text


def extract_text_from_pdf(file):
    file_stream = io.BytesIO(file.read())
    pdf_does_have_text = pdf_has_text(file_stream)
    if not pdf_does_have_text:
        print(f"extract_text_from_pdf pdf_has_text False")
        print(f"extract_text_from_pdf {str(file_stream)}")
        extracted_text_with_layout = perform_general_ocr(file)
        # extracted_text_with_col_line_num = process_for_line_column_numbers_with_desc_detection(extracted_text_with_layout)
        # print(f"extract_text_from_pdf {extracted_text_with_col_line_num[:500]}")
        return extracted_text_with_layout
    print(f"extract_text_from_pdf pdf_has_text True")
    pdf = PdfReader(file)
    text = "".join(page.extract_text() for page in pdf.pages)
    return text


def add_http_prefix(url):
    if url.startswith("http://") or url.startswith("https://"):
        pass
    else:
        url = "http://" + url
    return url


def hash_file_or_string(data):
    hasher = hashlib.sha256()
    if isinstance(data, str):
        hasher.update(data.encode('utf-8'))
    else:
        chunk = data.read(4096)
        while len(chunk) > 0:
            # Update the hash object with the current chunk
            hasher.update(chunk)
            chunk = data.read(4096)
    return hasher.hexdigest()


def is_image(file):
    mime_type = file.mimetype
    return mime_type and mime_type.startswith('image')


def encode_image_to_base64(file):
    file.seek(0)
    return base64.b64encode(file.read()).decode('utf-8')
