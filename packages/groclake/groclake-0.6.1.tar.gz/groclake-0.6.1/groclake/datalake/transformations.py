import fitz  # PyMuPDF
import base64
import magic
import io
import docx
import csv
from pptx import Presentation
import pandas as pd
import requests
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4

class Transformations:
    @staticmethod
    def chunk_text_manual(text, chunk_size=1000, overlap=200):
        """
        Splits text into overlapping chunks.
        :param text: The input string to chunk.
        :param chunk_size: The size of each chunk.
        :param overlap: The number of overlapping characters between chunks.
        :return: List of text chunks.
        """
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start += chunk_size - overlap
        return chunks

    @staticmethod
    def extract_text_from_pdf(file_bytes):
        """
        Extracts text from PDF bytes using PyMuPDF (fitz).
        :param file_bytes: PDF file content as bytes.
        :return: Extracted text as a string.
        """
        text = ""
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()
        return text.strip()
    
    @staticmethod
    def clean_up_text(text):
        """
        Clean large text files for easier processing while preserving context.
        
        Args:
            text (str): The input text to be cleaned
            
        Returns:
            str: Cleaned text
        """
        import re
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Fix spacing around punctuation
        text = re.sub(r'\s([.,;:!?)])', r'\1', text)
        text = re.sub(r'([({])\s', r'\1', text)
        
        # Normalize quotes
        text = re.sub(r'["""]', '"', text)
        
        # Fix common encoding issues
        text = re.sub(r'â€™', "'", text)
        text = re.sub(r'â€œ|â€', '"', text)
        text = re.sub(r'â€¦', '...', text)
        text = re.sub(r'â€"', '-', text)
        
        # Remove control characters while preserving newlines
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        
        # Normalize newlines (keep paragraph structure)
        text = re.sub(r'\n+', '\n', text)
        
        # Trim leading/trailing whitespace
        text = text.strip()
        
        # Remove any HTML tags that might be present
        text = re.sub(r'<[^>]+>', '', text)
        
        # Fix spacing after periods, question marks, and exclamation points
        text = re.sub(r'([.!?])\s*([A-Za-z])', r'\1 \2', text)
        
        # Remove URLs
        text = re.sub(r'https?://\S+', '', text)
        
        return text
    
    @staticmethod
    def get_file_type_from_base64(base64_string):
        decoded_data = base64.b64decode(base64_string)
        mime = magic.from_buffer(decoded_data, mime=True)
        if mime == 'application/pdf':
            return 'pdf'
        elif mime == 'text/plain':
            return 'txt'
        elif mime == 'image/jpeg':
            return 'jpg'
        elif mime == 'image/png':
            return 'png'
        elif mime == 'image/gif':
            return 'gif'
        elif mime == 'image/webp':
            return 'webp'
        elif mime == 'image/tiff':
            return 'tiff'
        elif mime == 'image/bmp':
            return 'bmp'
        elif mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return 'docx'
        elif mime == 'application/vnd.ms-excel':
            return 'xls'
        elif mime == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return 'xlsx'
        elif mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return 'pptx'
        elif mime == 'application/vnd.ms-powerpoint':
            return 'ppt'
        elif mime == 'application/csv':
            return 'csv'
        elif mime == 'audio/mpeg':
            return 'mp3'
        elif mime == 'audio/mp3':
            return 'mp3'
        elif mime == 'audio/mp4':
            return 'm4a'
            

        return mime
    
    @staticmethod
    def get_file_category_from_file_type(file_type: str):
        
        file_type = file_type.lower()
        
        if file_type == 'pdf':
            return 'pdf'
        elif file_type == 'txt':
            return 'text'
        elif file_type == 'jpg':
            return 'image'
        elif file_type == 'png':
            return 'image'
        elif file_type == 'gif':
            return 'image'
        elif file_type == 'webp':
            return 'image'
        elif file_type == 'tiff':
            return 'image'
        elif file_type == 'bmp':
            return 'image'
        elif file_type == 'docx':
            return 'word'
        elif file_type == 'xls':
            return 'excel'
        elif file_type == 'xlsx':
            return 'excel'
        elif file_type == 'pptx':
            return 'ppt'
        elif file_type == 'ppt':
            return 'ppt'
        elif file_type == 'csv':
            return 'csv'
        elif file_type == 'mp3':
            return 'audio'
        elif file_type == 'm4a':
            return 'audio'
        elif file_type == 'mp4':
            return 'video'
        elif file_type == 'm4v':
            return 'video'
        elif file_type == 'mov':
            return 'video'
        elif file_type == 'avi':
            return 'video'
        elif file_type == 'wmv':
            return 'video'
        elif file_type == 'js':
            return 'javascript'
        elif file_type == 'css':
            return 'css'
        elif file_type == 'html':
            return 'html'
        elif file_type == 'xml':
            return 'xml'
        elif file_type == 'json':
            return 'json'
        elif file_type == 'yaml':
            return 'yaml'
        elif file_type == 'yml':
            return 'yaml'
        elif file_type == 'toml':
            return 'toml'
        elif file_type == 'audio':
            return 'audio'
        elif file_type == 'video':
            return 'video'
        elif file_type == 'image':
            return 'image'
        elif file_type == 'text':
            return 'text'
        elif file_type == 'word':
            return 'word'
        elif file_type == 'excel':
            return 'excel'
        elif file_type == 'ppt':
            return 'ppt'
        
        return 'other'
    
    @staticmethod
    def get_file_type_from_bytes(bytes_data):
        mime = magic.from_buffer(bytes_data, mime=True)
        if mime == 'application/pdf':
            return 'pdf'
        elif mime == 'text/plain':
            return 'txt'
        elif mime == 'image/jpeg':
            return 'jpg'
        elif mime == 'image/png':
            return 'png'
        elif mime == 'image/gif':
            return 'gif'
        elif mime == 'image/webp':
            return 'webp'
        elif mime == 'image/tiff':
            return 'tiff'
        elif mime == 'image/bmp':
            return 'bmp'
        elif mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return 'docx'
        elif mime == 'application/vnd.ms-excel':
            return 'xls'
        elif mime == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return 'xlsx'
        elif mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return 'pptx'
        elif mime == 'application/vnd.ms-powerpoint':
            return 'ppt'
        elif mime == 'application/csv':
            return 'csv'
        elif mime == 'audio/mpeg':
            return 'mp3'
        elif mime == 'audio/mp3':
            return 'mp3'
        elif mime == 'audio/mp4':
            return 'm4a'

        return mime
    
    @staticmethod
    def transform_doc_to_text(document_bytes, file_type):
        """
        Transforms a document to text using the appropriate library based on the file type.
        
        Args:
            document_bytes (bytes): The bytes of the document to transform
        Returns:
            tuple: A tuple containing the text and page count
        """
        raw_text = ""
        page_count = 1
        if not file_type:
            file_type = Transformations.get_file_type_from_bytes(document_bytes)

        if file_type == 'pdf':
                pdf_stream = io.BytesIO(document_bytes)
                # Open with PyMuPDF
                doc = fitz.open(stream=pdf_stream, filetype="pdf")
                # Get page count
                page_count = len(doc)
                # Extract text (limit for summarization, optional chunking needed for long files)
                raw_text = "\n".join([page.get_text() for page in doc])
            
        elif file_type == 'docx':
            doc = docx.Document(io.BytesIO(document_bytes))  # Wrap bytes with BytesIO
            raw_text = "\n".join([para.text for para in doc.paragraphs])
            page_count = 1  # Page count not available in docx
        
        elif file_type == 'txt':
            raw_text = document_bytes.decode('utf-8')
            page_count = 1
        elif file_type == 'csv':
            text_lines = []
            with io.StringIO(document_bytes.decode('utf-8')) as csvfile:
                reader = csv.reader(csvfile)
                for row in reader:
                    text_lines.append(", ".join(row))
            raw_text = "\n".join(text_lines)
            page_count = 1
        elif file_type in ['xlsx', 'xls']:
            # Wrap raw bytes in BytesIO for pandas
            excel_io = io.BytesIO(document_bytes)
            # Engine depends on file type
            engine = 'openpyxl' if file_type == 'xlsx' else 'xlrd'
            # Read all sheets
            dfs = pd.read_excel(excel_io, sheet_name=None, engine=engine)
            # Combine all sheets into text (optional)
            raw_text = ""
            for sheet_name, df in dfs.items():
                raw_text += f"Sheet: {sheet_name}\n"
                raw_text += df.to_string(index=False)
                raw_text += "\n\n"

            page_count = len(dfs)

        elif file_type == 'pptx':
            prs = Presentation(io.BytesIO(document_bytes))  # Open from memory
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)

            raw_text = "\n".join(text_runs)
            page_count = len(prs.slides)  # Slide count
        elif file_type == 'ppt':
            raise ValueError("Legacy .ppt format is not supported. Please convert to .pptx.")
        elif file_type == 'jpg':
            raw_text = document_bytes.decode('utf-8')
            page_count = 1

        raw_text = Transformations.clean_up_text(raw_text)

        text_payload = {
            "raw_text": raw_text,
            "page_count": page_count
        }
        
        return text_payload

    @staticmethod
    def url_pdf_extract_template(pdf_url):

        response = requests.get(pdf_url)
        response.raise_for_status()  # Raise an error for HTTP issues

        doc = fitz.open(stream=BytesIO(response.content), filetype="pdf")
        pages = []

        for page in doc:
            blocks = page.get_text("dict")["blocks"]
            content = ""
            for block in blocks:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        content += span["text"] + " "
                content += "\n\n"
            pages.append(content.strip())

        return pages
    
    @staticmethod
    def create_pdf_in_buffer(text):
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=A4)
        width, height = A4
        y = height - 50  # start near top

        for line in text.split('\n'):
            c.drawString(50, y, line)
            y -= 15
            if y < 50:  # new page if needed
                c.showPage()
                y = height - 50

        c.save()
        buffer.seek(0)  # Important: rewind the buffer before returning
        return buffer

